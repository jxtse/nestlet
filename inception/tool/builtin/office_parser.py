"""
Office document parsing tools.

Provides tools for parsing Microsoft Office documents:
- Word (.docx)
- Excel (.xlsx)
- PowerPoint (.pptx)
"""

from __future__ import annotations

import logging
import time
from pathlib import Path
from typing import Any, Dict, List, Optional

from inception.tool.base import (
    Tool,
    ToolSpec,
    ToolResult,
    ParameterSpec,
    ParameterType,
    ReturnSpec,
)

logger = logging.getLogger(__name__)


class ParseWordTool(Tool):
    """Tool for parsing Word documents (.docx)."""

    def __init__(self):
        self._spec = ToolSpec(
            name="parse_word",
            description="Parse a Microsoft Word document (.docx) and extract its text content, including paragraphs, tables, and headers.",
            parameters={
                "file_path": ParameterSpec(
                    name="file_path",
                    type=ParameterType.STRING,
                    description="Path to the Word document (.docx file)",
                    required=True,
                ),
                "include_tables": ParameterSpec(
                    name="include_tables",
                    type=ParameterType.BOOLEAN,
                    description="Whether to include table contents (default: True)",
                    required=False,
                    default=True,
                ),
            },
            returns=ReturnSpec(
                type=ParameterType.STRING,
                description="Extracted text content from the Word document",
            ),
            category="office",
            tags=["word", "docx", "document", "parser"],
        )

    @property
    def spec(self) -> ToolSpec:
        return self._spec

    async def execute(self, **kwargs: Any) -> ToolResult:
        start_time = time.time()
        file_path = kwargs.get("file_path", "")
        include_tables = kwargs.get("include_tables", True)

        if not file_path:
            return ToolResult.fail(
                error="file_path is required",
                execution_time=time.time() - start_time,
            )

        path = Path(file_path)
        if not path.exists():
            return ToolResult.fail(
                error=f"File not found: {file_path}",
                execution_time=time.time() - start_time,
            )

        if path.suffix.lower() != ".docx":
            return ToolResult.fail(
                error=f"Invalid file type: {path.suffix}. Expected .docx",
                execution_time=time.time() - start_time,
            )

        try:
            from docx import Document

            doc = Document(file_path)
            content_parts = []

            # Extract paragraphs
            for para in doc.paragraphs:
                text = para.text.strip()
                if text:
                    # Check if it's a heading
                    if para.style.name.startswith("Heading"):
                        level = para.style.name.replace("Heading ", "")
                        content_parts.append(f"\n{'#' * int(level) if level.isdigit() else '##'} {text}\n")
                    else:
                        content_parts.append(text)

            # Extract tables
            if include_tables and doc.tables:
                for i, table in enumerate(doc.tables):
                    content_parts.append(f"\n[Table {i + 1}]")
                    for row in table.rows:
                        row_text = " | ".join(cell.text.strip() for cell in row.cells)
                        content_parts.append(f"| {row_text} |")

            result = "\n".join(content_parts)
            return ToolResult.ok(
                result=result,
                execution_time=time.time() - start_time,
            )

        except ImportError:
            return ToolResult.fail(
                error="python-docx is not installed. Install with: pip install python-docx",
                execution_time=time.time() - start_time,
            )
        except Exception as e:
            return ToolResult.fail(
                error=f"Failed to parse Word document: {str(e)}",
                execution_time=time.time() - start_time,
            )


class ParseExcelTool(Tool):
    """Tool for parsing Excel spreadsheets (.xlsx and .xls)."""

    def __init__(self):
        self._spec = ToolSpec(
            name="parse_excel",
            description="Parse a Microsoft Excel spreadsheet (.xlsx or .xls) and extract data from sheets. Supports both new (.xlsx) and legacy (.xls) Excel formats. Returns data as formatted text or JSON.",
            parameters={
                "file_path": ParameterSpec(
                    name="file_path",
                    type=ParameterType.STRING,
                    description="Path to the Excel file (.xlsx or .xls)",
                    required=True,
                ),
                "sheet_name": ParameterSpec(
                    name="sheet_name",
                    type=ParameterType.STRING,
                    description="Name of the sheet to parse. If not specified, parses all sheets.",
                    required=False,
                    default=None,
                ),
                "max_rows": ParameterSpec(
                    name="max_rows",
                    type=ParameterType.INTEGER,
                    description="Maximum number of rows to read (default: 1000)",
                    required=False,
                    default=1000,
                ),
                "output_format": ParameterSpec(
                    name="output_format",
                    type=ParameterType.STRING,
                    description="Output format: 'text' (markdown table) or 'json' (default: 'text')",
                    required=False,
                    default="text",
                ),
            },
            returns=ReturnSpec(
                type=ParameterType.STRING,
                description="Extracted data from the Excel file",
            ),
            category="office",
            tags=["excel", "xlsx", "xls", "spreadsheet", "parser"],
        )

    @property
    def spec(self) -> ToolSpec:
        return self._spec

    def _detect_excel_format(self, file_path: str) -> str:
        """Detect the actual Excel format by reading file header."""
        with open(file_path, 'rb') as f:
            header = f.read(8)

        # Check for ZIP signature (xlsx format)
        if header[:4] == b'PK\x03\x04':
            return "xlsx"
        # Check for OLE2 signature (xls format)
        elif header[:8] == b'\xd0\xcf\x11\xe0\xa1\xb1\x1a\xe1':
            return "xls"
        else:
            return "unknown"

    async def execute(self, **kwargs: Any) -> ToolResult:
        start_time = time.time()
        file_path = kwargs.get("file_path", "")
        sheet_name = kwargs.get("sheet_name")
        max_rows = kwargs.get("max_rows", 1000)
        output_format = kwargs.get("output_format", "text")

        if not file_path:
            return ToolResult.fail(
                error="file_path is required",
                execution_time=time.time() - start_time,
            )

        path = Path(file_path)
        if not path.exists():
            return ToolResult.fail(
                error=f"File not found: {file_path}",
                execution_time=time.time() - start_time,
            )

        if path.suffix.lower() not in (".xlsx", ".xls"):
            return ToolResult.fail(
                error=f"Invalid file type: {path.suffix}. Expected .xlsx or .xls",
                execution_time=time.time() - start_time,
            )

        # Detect actual file format (extension might not match actual format)
        actual_format = self._detect_excel_format(file_path)

        try:
            import json

            if actual_format == "xls":
                # Use xlrd for legacy OLE2 format
                return await self._parse_xls(file_path, sheet_name, max_rows, output_format, start_time)
            elif actual_format == "xlsx":
                # Use openpyxl for modern xlsx format
                return await self._parse_xlsx(file_path, sheet_name, max_rows, output_format, start_time)
            else:
                return ToolResult.fail(
                    error=f"Unknown Excel format. File header does not match xlsx or xls format.",
                    execution_time=time.time() - start_time,
                )

        except ImportError as e:
            return ToolResult.fail(
                error=f"Missing dependency: {str(e)}. Install with: uv pip install openpyxl xlrd",
                execution_time=time.time() - start_time,
            )
        except Exception as e:
            return ToolResult.fail(
                error=f"Failed to parse Excel file: {str(e)}",
                execution_time=time.time() - start_time,
            )

    async def _parse_xls(self, file_path: str, sheet_name: Optional[str], max_rows: int, output_format: str, start_time: float) -> ToolResult:
        """Parse legacy .xls (OLE2) format using xlrd."""
        try:
            import xlrd
            import json

            wb = xlrd.open_workbook(file_path)
            result_parts = []

            sheets_to_process = [sheet_name] if sheet_name else wb.sheet_names()

            for sname in sheets_to_process:
                if sname not in wb.sheet_names():
                    continue

                ws = wb.sheet_by_name(sname)
                result_parts.append(f"\n## Sheet: {sname}\n")

                rows_data = []
                for i in range(min(ws.nrows, max_rows)):
                    row = [ws.cell_value(i, j) for j in range(ws.ncols)]
                    rows_data.append(row)

                if ws.nrows > max_rows:
                    result_parts.append(f"\n... (truncated at {max_rows} rows, total: {ws.nrows})")

                if not rows_data:
                    result_parts.append("(empty sheet)")
                    continue

                if output_format == "json":
                    headers = [str(h) if h else f"col_{i}" for i, h in enumerate(rows_data[0])]
                    data = []
                    for row in rows_data[1:]:
                        row_dict = {}
                        for j, val in enumerate(row):
                            if j < len(headers):
                                row_dict[headers[j]] = val
                        data.append(row_dict)
                    result_parts.append(json.dumps(data, indent=2, default=str, ensure_ascii=False))
                else:
                    for i, row in enumerate(rows_data):
                        row_str = " | ".join(str(cell) if cell is not None else "" for cell in row)
                        result_parts.append(f"| {row_str} |")
                        if i == 0:
                            separator = " | ".join("---" for _ in row)
                            result_parts.append(f"| {separator} |")

            result = "\n".join(result_parts)
            return ToolResult.ok(
                result=result,
                execution_time=time.time() - start_time,
            )

        except ImportError:
            return ToolResult.fail(
                error="xlrd is not installed. Install with: uv pip install xlrd",
                execution_time=time.time() - start_time,
            )

    async def _parse_xlsx(self, file_path: str, sheet_name: Optional[str], max_rows: int, output_format: str, start_time: float) -> ToolResult:
        """Parse modern .xlsx format using openpyxl."""
        from openpyxl import load_workbook
        import json

        wb = load_workbook(file_path, read_only=True, data_only=True)
        result_parts = []

        sheets_to_process = [sheet_name] if sheet_name else wb.sheetnames

        for sname in sheets_to_process:
            if sname not in wb.sheetnames:
                continue

            ws = wb[sname]
            result_parts.append(f"\n## Sheet: {sname}\n")

            rows_data = []
            for i, row in enumerate(ws.iter_rows(values_only=True)):
                if i >= max_rows:
                    result_parts.append(f"\n... (truncated at {max_rows} rows)")
                    break
                rows_data.append(row)

            if not rows_data:
                result_parts.append("(empty sheet)")
                continue

            if output_format == "json":
                # Use first row as headers
                headers = [str(h) if h else f"col_{i}" for i, h in enumerate(rows_data[0])]
                data = []
                for row in rows_data[1:]:
                    row_dict = {}
                    for j, val in enumerate(row):
                        if j < len(headers):
                            row_dict[headers[j]] = val
                    data.append(row_dict)
                result_parts.append(json.dumps(data, indent=2, default=str, ensure_ascii=False))
            else:
                # Markdown table format
                for i, row in enumerate(rows_data):
                    row_str = " | ".join(str(cell) if cell is not None else "" for cell in row)
                    result_parts.append(f"| {row_str} |")
                    if i == 0:
                        # Add separator after header
                        separator = " | ".join("---" for _ in row)
                        result_parts.append(f"| {separator} |")

        wb.close()
        result = "\n".join(result_parts)
        return ToolResult.ok(
            result=result,
            execution_time=time.time() - start_time,
        )


class ParsePowerPointTool(Tool):
    """Tool for parsing PowerPoint presentations (.pptx)."""

    def __init__(self):
        self._spec = ToolSpec(
            name="parse_powerpoint",
            description="Parse a Microsoft PowerPoint presentation (.pptx) and extract text content from slides, including titles, body text, and notes.",
            parameters={
                "file_path": ParameterSpec(
                    name="file_path",
                    type=ParameterType.STRING,
                    description="Path to the PowerPoint file (.pptx)",
                    required=True,
                ),
                "include_notes": ParameterSpec(
                    name="include_notes",
                    type=ParameterType.BOOLEAN,
                    description="Whether to include speaker notes (default: True)",
                    required=False,
                    default=True,
                ),
            },
            returns=ReturnSpec(
                type=ParameterType.STRING,
                description="Extracted text content from the PowerPoint presentation",
            ),
            category="office",
            tags=["powerpoint", "pptx", "presentation", "parser"],
        )

    @property
    def spec(self) -> ToolSpec:
        return self._spec

    async def execute(self, **kwargs: Any) -> ToolResult:
        start_time = time.time()
        file_path = kwargs.get("file_path", "")
        include_notes = kwargs.get("include_notes", True)

        if not file_path:
            return ToolResult.fail(
                error="file_path is required",
                execution_time=time.time() - start_time,
            )

        path = Path(file_path)
        if not path.exists():
            return ToolResult.fail(
                error=f"File not found: {file_path}",
                execution_time=time.time() - start_time,
            )

        if path.suffix.lower() != ".pptx":
            return ToolResult.fail(
                error=f"Invalid file type: {path.suffix}. Expected .pptx",
                execution_time=time.time() - start_time,
            )

        try:
            from pptx import Presentation

            prs = Presentation(file_path)
            content_parts = []

            for i, slide in enumerate(prs.slides, 1):
                content_parts.append(f"\n## Slide {i}\n")

                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        # Check if it's a title
                        if shape.is_placeholder and hasattr(shape, "placeholder_format"):
                            if shape.placeholder_format.type == 1:  # Title
                                content_parts.append(f"### {shape.text.strip()}\n")
                            else:
                                content_parts.append(shape.text.strip())
                        else:
                            content_parts.append(shape.text.strip())

                    # Handle tables in shapes
                    if shape.has_table:
                        table = shape.table
                        content_parts.append("\n[Table]")
                        for row in table.rows:
                            row_text = " | ".join(cell.text.strip() for cell in row.cells)
                            content_parts.append(f"| {row_text} |")

                # Extract notes
                if include_notes and slide.has_notes_slide:
                    notes_frame = slide.notes_slide.notes_text_frame
                    if notes_frame and notes_frame.text.strip():
                        content_parts.append(f"\n**Speaker Notes:** {notes_frame.text.strip()}")

            result = "\n".join(content_parts)
            return ToolResult.ok(
                result=result,
                execution_time=time.time() - start_time,
            )

        except ImportError:
            return ToolResult.fail(
                error="python-pptx is not installed. Install with: pip install python-pptx",
                execution_time=time.time() - start_time,
            )
        except Exception as e:
            return ToolResult.fail(
                error=f"Failed to parse PowerPoint file: {str(e)}",
                execution_time=time.time() - start_time,
            )


class ParsePDFTool(Tool):
    """Tool for parsing PDF documents."""

    def __init__(self):
        self._spec = ToolSpec(
            name="parse_pdf",
            description="Parse a PDF document and extract its text content.",
            parameters={
                "file_path": ParameterSpec(
                    name="file_path",
                    type=ParameterType.STRING,
                    description="Path to the PDF file",
                    required=True,
                ),
                "max_pages": ParameterSpec(
                    name="max_pages",
                    type=ParameterType.INTEGER,
                    description="Maximum number of pages to extract (default: 50)",
                    required=False,
                    default=50,
                ),
            },
            returns=ReturnSpec(
                type=ParameterType.STRING,
                description="Extracted text content from the PDF",
            ),
            category="office",
            tags=["pdf", "document", "parser"],
        )

    @property
    def spec(self) -> ToolSpec:
        return self._spec

    async def execute(self, **kwargs: Any) -> ToolResult:
        start_time = time.time()
        file_path = kwargs.get("file_path", "")
        max_pages = kwargs.get("max_pages", 50)

        if not file_path:
            return ToolResult.fail(
                error="file_path is required",
                execution_time=time.time() - start_time,
            )

        path = Path(file_path)
        if not path.exists():
            return ToolResult.fail(
                error=f"File not found: {file_path}",
                execution_time=time.time() - start_time,
            )

        if path.suffix.lower() != ".pdf":
            return ToolResult.fail(
                error=f"Invalid file type: {path.suffix}. Expected .pdf",
                execution_time=time.time() - start_time,
            )

        try:
            import pypdf

            reader = pypdf.PdfReader(file_path)
            content_parts = []

            total_pages = len(reader.pages)
            pages_to_read = min(total_pages, max_pages)

            content_parts.append(f"PDF Document: {path.name}")
            content_parts.append(f"Total pages: {total_pages}")
            if pages_to_read < total_pages:
                content_parts.append(f"(Reading first {pages_to_read} pages)")
            content_parts.append("")

            for i in range(pages_to_read):
                page = reader.pages[i]
                text = page.extract_text()
                if text.strip():
                    content_parts.append(f"\n--- Page {i + 1} ---\n")
                    content_parts.append(text.strip())

            result = "\n".join(content_parts)
            return ToolResult.ok(
                result=result,
                execution_time=time.time() - start_time,
            )

        except ImportError:
            return ToolResult.fail(
                error="pypdf is not installed. Install with: pip install pypdf",
                execution_time=time.time() - start_time,
            )
        except Exception as e:
            return ToolResult.fail(
                error=f"Failed to parse PDF file: {str(e)}",
                execution_time=time.time() - start_time,
            )
